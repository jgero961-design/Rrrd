#!/usr/bin/env python3
import math
import os
from dataclasses import dataclass
from typing import List, Tuple

import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_AUTO_SIZE
import xlsxwriter


OUTPUT_DIR = os.path.dirname(__file__)


@dataclass
class SiteSpec:
    land_length_m: float = 58.0
    land_width_m: float = 19.0
    driveway_width_m: float = 4.0
    num_duplexes_per_row: int = 6


@dataclass
class FinancialInputs:
    num_units: int = 12
    adr_gel: float = 150.0
    occupancy_rate: float = 0.75
    opex_pct_of_gross: float = 0.30
    mgmt_pct_of_gross: float = 0.10
    capex_low_gel: float = 385000.0
    capex_high_gel: float = 728000.0


def generate_site_layout_images(output_dir: str, spec: SiteSpec) -> Tuple[str, str]:
    """Generate site layout PNG and SVG with basic dimension annotations."""
    png_path = os.path.join(output_dir, "site_layout.png")
    svg_path = os.path.join(output_dir, "site_layout.svg")

    fig_width_in = spec.land_length_m * 0.12  # scale factor for display
    fig_height_in = spec.land_width_m * 0.12
    fig, ax = plt.subplots(figsize=(fig_width_in, fig_height_in), dpi=200)

    # Coordinate system: x in [0, length], y in [0, width]
    ax.set_xlim(0, spec.land_length_m)
    ax.set_ylim(0, spec.land_width_m)
    ax.set_aspect('equal')
    ax.axis('off')

    # Plot boundary
    plot_rect = Rectangle((0, 0), spec.land_length_m, spec.land_width_m,
                          linewidth=2, edgecolor='black', facecolor='none')
    ax.add_patch(plot_rect)

    # Driveway centered across the length
    driveway_y0 = (spec.land_width_m - spec.driveway_width_m) / 2.0
    driveway_rect = Rectangle((0, driveway_y0), spec.land_length_m, spec.driveway_width_m,
                              linewidth=1.2, edgecolor='gray', facecolor='#e6e6e6', alpha=0.9)
    ax.add_patch(driveway_rect)

    # Rows above and below driveway
    row_height = (spec.land_width_m - spec.driveway_width_m) / 2.0

    # Draw units as bays (note: conceptual bay sizing to fit plot; not to container depth)
    bays = spec.num_duplexes_per_row
    bay_len = spec.land_length_m / bays
    margin = bay_len * 0.07

    # Top row (odd units)
    top_y0 = driveway_y0 + spec.driveway_width_m
    top_y1 = spec.land_width_m
    # Bottom row (even units)
    bot_y0 = 0
    bot_y1 = driveway_y0

    def draw_row(is_top: bool):
        y0, y1 = (top_y0, top_y1) if is_top else (bot_y0, bot_y1)
        unit_label_start = 1 if is_top else 2
        yard_depth = (y1 - y0) * 0.45
        garage_depth = (y1 - y0) * 0.35
        gap_depth = (y1 - y0) - (yard_depth + garage_depth)
        # yard near outer edge; garage near driveway
        if is_top:
            yard_rect = lambda x0, x1: (x0 + margin, y1 - yard_depth - 0.02*(y1-y0), x1 - x0 - 2*margin, yard_depth)
            garage_rect = lambda x0, x1: (x0 + margin, y0 + 0.02*(y1-y0), x1 - x0 - 2*margin, garage_depth)
            unit_text_y = y0 + garage_depth + gap_depth/2.0
        else:
            yard_rect = lambda x0, x1: (x0 + margin, y0 + 0.02*(y1-y0), x1 - x0 - 2*margin, yard_depth)
            garage_rect = lambda x0, x1: (x0 + margin, y1 - garage_depth - 0.02*(y1-y0), x1 - x0 - 2*margin, garage_depth)
            unit_text_y = y1 - garage_depth - gap_depth/2.0

        for i in range(bays):
            x0 = i * bay_len
            x1 = (i + 1) * bay_len
            # Yard
            rx, ry, rw, rh = yard_rect(x0, x1)
            ax.add_patch(Rectangle((rx, ry), rw, rh, edgecolor='#4CAF50', facecolor='#dff0d8', linewidth=1.0))
            # Garage/Living block
            rx2, ry2, rw2, rh2 = garage_rect(x0, x1)
            ax.add_patch(Rectangle((rx2, ry2), rw2, rh2, edgecolor='#1976D2', facecolor='#e3f2fd', linewidth=1.0))

            unit_number = unit_label_start + i * 2
            ax.text((x0 + x1)/2.0, unit_text_y, f"Unit {unit_number}", fontsize=7, ha='center', va='center', color='black')

    draw_row(is_top=True)
    draw_row(is_top=False)

    # Labels
    ax.text(spec.land_length_m/2.0, driveway_y0 + spec.driveway_width_m/2.0, "Central Driveway\n~4 m",
            fontsize=8, ha='center', va='center', color='dimgray')

    # Dimension annotations
    # Horizontal (length) dimension
    ax.add_line(Line2D([0, spec.land_length_m], [-0.6, -0.6], color='black', linewidth=0.8))
    ax.add_line(Line2D([0, 0], [-0.4, -0.8], color='black', linewidth=0.8))
    ax.add_line(Line2D([spec.land_length_m, spec.land_length_m], [-0.4, -0.8], color='black', linewidth=0.8))
    ax.text(spec.land_length_m/2.0, -1.3, f"{spec.land_length_m:.0f} m", fontsize=8, ha='center', va='top')

    # Vertical (width) dimension
    ax.add_line(Line2D([-0.6, -0.6], [0, spec.land_width_m], color='black', linewidth=0.8))
    ax.add_line(Line2D([-0.4, -0.8], [0, 0], color='black', linewidth=0.8))
    ax.add_line(Line2D([-0.4, -0.8], [spec.land_width_m, spec.land_width_m], color='black', linewidth=0.8))
    ax.text(-1.2, spec.land_width_m/2.0, f"{spec.land_width_m:.0f} m", fontsize=8, ha='center', va='center', rotation=90)

    # Title
    ax.text(spec.land_length_m/2.0, spec.land_width_m + 0.8,
            "Container House Project — Site Layout (Not to container depth; scaled to plot)",
            fontsize=10, ha='center', va='bottom', weight='bold')

    plt.tight_layout(pad=1.2)
    fig.savefig(png_path, bbox_inches='tight')
    fig.savefig(svg_path, bbox_inches='tight')
    plt.close(fig)

    return png_path, svg_path


def generate_financial_model(output_dir: str, inputs: FinancialInputs) -> str:
    """Generate an Excel pro-forma with inputs, summary, monthly projection, sensitivity, and cap table."""
    xlsx_path = os.path.join(output_dir, "financial_model.xlsx")
    wb = xlsxwriter.Workbook(xlsx_path)

    # Formats
    fmt_bold = wb.add_format({'bold': True})
    fmt_money = wb.add_format({'num_format': '#,##0 [$₾-420]'} )
    fmt_pct = wb.add_format({'num_format': '0%'} )
    fmt_title = wb.add_format({'bold': True, 'font_size': 14})
    fmt_subtitle = wb.add_format({'bold': True, 'font_size': 11})
    fmt_int = wb.add_format({'num_format': '#,##0'})
    fmt_text = wb.add_format({})

    # Sheet: Inputs
    ws_in = wb.add_worksheet('Inputs')
    ws_in.set_column('A:A', 36)
    ws_in.set_column('B:B', 18)
    ws_in.write('A1', 'Financial Model Inputs', fmt_title)

    inputs_rows = [
        ('Number of units', inputs.num_units, fmt_int),
        ('Average Daily Rate (ADR) — GEL', inputs.adr_gel, fmt_money),
        ('Occupancy rate', inputs.occupancy_rate, fmt_pct),
        ('Operating expenses (% of gross)', inputs.opex_pct_of_gross, fmt_pct),
        ('Management fee (% of gross)', inputs.mgmt_pct_of_gross, fmt_pct),
        ('CAPEX (low) — GEL', inputs.capex_low_gel, fmt_money),
        ('CAPEX (high) — GEL', inputs.capex_high_gel, fmt_money),
    ]

    for idx, (label, value, value_fmt) in enumerate(inputs_rows, start=3):
        ws_in.write(f'A{idx}', label)
        ws_in.write_number(f'B{idx}', value, value_fmt)

    # Sheet: Summary
    ws = wb.add_worksheet('Summary')
    ws.set_column('A:A', 40)
    ws.set_column('B:B', 20)
    ws.write('A1', 'Summary Outputs', fmt_title)

    ws.write('A3', 'Nights per unit per year', fmt_bold)
    ws.write_number('B3', 365, fmt_int)

    ws.write('A4', 'Booked nights per unit', fmt_bold)
    ws.write_formula('B4', '=B3*Inputs!B4', fmt_int)  # 365 * occupancy

    ws.write('A5', 'Revenue per unit (GEL)', fmt_bold)
    ws.write_formula('B5', '=Inputs!B3*B4', fmt_money)  # ADR * booked nights per unit

    ws.write('A6', 'Total gross revenue (GEL)', fmt_bold)
    ws.write_formula('B6', '=Inputs!B2*B5', fmt_money)  # num_units * revenue per unit

    ws.write('A8', 'Operating expenses (GEL)', fmt_bold)
    ws.write_formula('B8', '=Inputs!B5*Inputs!B2*Inputs!B5/Inputs!B3', fmt_money)  # placeholder to avoid Excel name clash

    # Correct OPEX and MGMT based on gross
    ws.write('A9', 'Operating expenses (GEL) — 30% of gross', fmt_bold)
    ws.write_formula('B9', '=B6*Inputs!B5', fmt_money)  # opex% at Inputs!B5

    ws.write('A10', 'Management fee (GEL) — 10% of gross', fmt_bold)
    ws.write_formula('B10', '=B6*Inputs!B6', fmt_money)

    ws.write('A12', 'Net distributable cash (GEL)', fmt_title)
    ws.write_formula('B12', '=B6-B9-B10', fmt_money)

    # CAPEX range
    ws.write('A14', 'CAPEX (low) — GEL', fmt_bold)
    ws.write_formula('B14', '=Inputs!B7', fmt_money)
    ws.write('A15', 'CAPEX (high) — GEL', fmt_bold)
    ws.write_formula('B15', '=Inputs!B8', fmt_money)

    # Implied annual ROI on CAPEX mid-point
    ws.write('A17', 'Implied annual ROI on CAPEX midpoint', fmt_bold)
    ws.write_formula('B17', '=B12/AVERAGE(B14,B15)', fmt_pct)

    # Sheet: Monthly Projection
    ws_m = wb.add_worksheet('MonthlyProjection')
    ws_m.set_column('A:A', 20)
    ws_m.set_column('B:D', 18)

    months = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
    ws_m.write('A1', 'Monthly Projection (equalized for demo)', fmt_subtitle)
    ws_m.write_row('A3', ['Month', 'Gross Revenue', 'Operating Expenses', 'Management Fee'])

    # Gross annual from Summary!B6
    ws_m.write('F2', 'Annual Gross (link)')
    ws_m.write_formula('G2', '=Summary!B6')
    ws_m.write('F3', 'Annual OPEX (link)')
    ws_m.write_formula('G3', '=Summary!B9')
    ws_m.write('F4', 'Annual Mgmt (link)')
    ws_m.write_formula('G4', '=Summary!B10')

    for i, m in enumerate(months, start=4):
        ws_m.write(f'A{i}', m)
        ws_m.write_formula(f'B{i}', '=G2/12', fmt_money)
        ws_m.write_formula(f'C{i}', '=G3/12', fmt_money)
        ws_m.write_formula(f'D{i}', '=G4/12', fmt_money)

    ws_m.write('A18', 'Monthly Net Distributable (per month)', fmt_bold)
    ws_m.write_formula('B18', '=Summary!B12/12', fmt_money)

    # Sheet: Sensitivity (ADR x Occupancy)
    ws_s = wb.add_worksheet('Sensitivity')
    ws_s.set_column('A:A', 18)
    ws_s.set_column('B:E', 16)
    ws_s.write('A1', 'Sensitivity: Net Distributable GEL vs ADR/Occupancy', fmt_subtitle)

    adr_values = [120, 150, 180, 210]
    occ_values = [0.60, 0.75, 0.90]

    ws_s.write_row('B2', [f'ADR ₾{v}' for v in adr_values], fmt_bold)
    for r, occ in enumerate(occ_values, start=3):
        ws_s.write(f'A{r}', f'Occ {int(occ*100)}%', fmt_bold)
        for c, adr in enumerate(adr_values, start=2):
            # Formula for net: num_units * ADR * 365 * occ * (1 - opex - mgmt)
            ws_s.write_formula(r-1, c-1, f"=Inputs!B2*{adr}*365*{occ}*(1-Inputs!B5-Inputs!B6)", fmt_money)

    # Sheet: Cap Table and Payout
    ws_c = wb.add_worksheet('CapTable')
    ws_c.set_column('A:A', 34)
    ws_c.set_column('B:C', 16)
    ws_c.write('A1', 'Cap Table & Payout Calculator', fmt_subtitle)

    ws_c.write('A3', 'Ownership % (enter)')
    ws_c.write('B3', 'Monthly Payout (GEL)')

    # Example rows for 1% to 10%
    for i in range(1, 11):
        row = 3 + i
        ws_c.write_number(row-1, 0, i/100, fmt_pct)
        ws_c.write_formula(row-1, 1, '=Summary!B12/12*A{}'.format(row), fmt_money)

    # Specific example: 1% monthly payout expected
    ws_c.write('A16', 'Example: 1% share monthly payout (excerpt)', fmt_bold)
    ws_c.write_formula('B16', '=Summary!B12/12*1%', fmt_money)

    wb.close()
    return xlsx_path


def _add_bullets(shape, lines: List[str]):
    tf = shape.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0


def generate_pitch_deck(output_dir: str) -> str:
    pptx_path = os.path.join(output_dir, 'pitch_deck.pptx')
    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # Slide 1: Project Overview & Vision
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = 'Container House Investment — Overview & Vision'
    s1.placeholders[1].text = 'Tbilisi, Georgia — 58m × 19m plot; 12 modular units with garages, yards, rooftop verandas'

    # Slide 2: Site Layout & Blueprint
    s2 = prs.slides.add_slide(bullet_layout)
    s2.shapes.title.text = 'Site Layout & Blueprint'
    _add_bullets(s2.shapes.placeholders[1], [
        'Two rows of 6 duplexes (12 units total) with central ~4 m driveway',
        'Each unit: garage, living, private yard, rooftop veranda',
        'Scaled site layout with dimensions in appendix'
    ])

    # Slide 3: Market & Financial Highlights
    s3 = prs.slides.add_slide(bullet_layout)
    s3.shapes.title.text = 'Market & Financial Highlights'
    _add_bullets(s3.shapes.placeholders[1], [
        'ADR: ₾150; Occupancy: 75% (conservative)',
        'Annual gross revenue: ₾492,750 (est.)',
        'Net distributable cash: ~₾271,000 after OPEX (30%) and management (10%)'
    ])

    # Slide 4: Ownership & Investment Model
    s4 = prs.slides.add_slide(bullet_layout)
    s4.shapes.title.text = 'Ownership & Investment Model'
    _add_bullets(s4.shapes.placeholders[1], [
        'Fractional ownership in SPV; distributions proportional to shares',
        'Minimum buy-in: 1% (~₾5,569)',
        'Operator retains 10–20% equity plus management fee'
    ])

    # Slide 5: Legal & Operational Plan
    s5 = prs.slides.add_slide(bullet_layout)
    s5.shapes.title.text = 'Legal & Operational Plan'
    _add_bullets(s5.shapes.placeholders[1], [
        'Form LLC/SPV; investor agreements; transfer/exit terms',
        'Permits for construction and short-term rentals',
        'Reserve funds; monthly reporting to investors'
    ])

    # Slide 6: Next Steps & Call to Action
    s6 = prs.slides.add_slide(bullet_layout)
    s6.shapes.title.text = 'Next Steps & Call to Action'
    _add_bullets(s6.shapes.placeholders[1], [
        'Finalize architectural drawings and site plans',
        'Build full financial pro-forma and sensitivity analysis',
        'Begin permitting; prepare for construction and investor onboarding'
    ])

    prs.save(pptx_path)
    return pptx_path


def generate_readme(output_dir: str) -> str:
    readme_path = os.path.join(output_dir, 'README.md')
    content = """
## Container House Investment Project — Full Documentation

1. Project Overview

- **Location**: Tbilisi, Georgia
- **Land Size**: 58 m (length) × 19 m (width)
- **Purpose**: Modular container housing units for daily rental with private yards, parking, and rooftop verandas
- **Units**: 12 total, arranged as 6 duplexes (2 units sharing a back wall), split in two rows with a central driveway
- **Ownership Model**: Fractional ownership in a single SPV with distributions proportional to share

2. Site Layout & Physical Blueprint

ASCII Diagram of Site Layout (Top-down view)

```
+---------------------------------------------------------+
| Duplex Row 1 (6 duplexes)                               |
| +--------+ +--------+ +--------+ +--------+ +--------+ +--------+ |
| | Unit 1 | | Unit 3 | | Unit 5 | | Unit 7 | | Unit 9 | | Unit 11| |
| | Garage | | Garage | | Garage | | Garage | | Garage | | Garage | |
| | Living | | Living | | Living | | Living | | Living | | Living | |
| | Yard   | | Yard   | | Yard   | | Yard   | | Yard   | | Yard   | |
| +--------+ +--------+ +--------+ +--------+ +--------+ +--------+ |
|                                                         |
|                  Central Driveway (~4 m wide)           |
|                                                         |
| +--------+ +--------+ +--------+ +--------+ +--------+ +--------+ |
| | Unit 2 | | Unit 4 | | Unit 6 | | Unit 8 | | Unit 10| | Unit 12| |
| | Garage | | Garage | | Garage | | Garage | | Garage | | Garage | |
| | Living | | Living | | Living | | Living | | Living | | Living | |
| | Yard   | | Yard   | | Yard   | | Yard   | | Yard   | | Yard   | |
| +--------+ +--------+ +--------+ +--------+ +--------+ +--------+ |
| Duplex Row 2 (6 duplexes)                               |
+---------------------------------------------------------+
```

- Each "Unit" is a container-based duplex side.
- The driveway gives access to all garages from the center.
- Each unit has a fenced yard including a firepit and grill spot.

Dimensions and Notes

- **Land length**: 58.0 m
- **Land width**: 19.0 m
- **Container length**: 12.19 m
- **Container width**: 2.44 m
- **Duplex length (container + yard)**: ~12.7 m
- **Duplex width (2 containers back-to-back)**: 4.88 m
- **Yard depth**: ~3.5 m
- **Central driveway width**: ~4.0 m
- **Total duplexes per row**: 6 (Total 12 units)

3. Structural Details per Unit

- **Ground floor**: Garage + storage for 1 car inside container footprint
- **Second floor**: Living space with kitchenette, bathroom, bedroom, stairs to veranda
- **Rooftop veranda**: Open-air leisure space with safety rails
- **Private fenced yard**: Space for firepit and grill, gate for car & pedestrian access

4. Financial Model Summary

- **Total units**: 12
- **Project CAPEX**: ₾385,000 - ₾728,000
- **Average Daily Rate (ADR)**: ₾150 / night
- **Occupancy**: 75%
- **Gross Annual Revenue**: ₾492,750 (based on ADR & occupancy)
- **Operating Expenses**: 30% of gross
- **Management Fee**: 10% of gross
- **Net Distributable Cash**: ~₾271,000 (after OPEX & fees)

5. Investment Structure & Returns

- Investors buy fractional ownership of the whole project via SPV
- **Minimum buy-in**: 1% (₾5,569 approx.)
- Profit distributed monthly proportional to ownership
- Operator retains 10–20% equity + management fee
- **Estimated annual ROI (after fees)**: ~48.7%
- **Example**: 1% share → ₾226/month in rental income payout

6. Legal & Operational Considerations

- Form LLC/SPV owning all assets
- Draft legal agreements covering ownership, distributions, transfers, and exit strategy
- Obtain building permits and short-term rental approvals
- Set reserve funds for maintenance & unexpected expenses
- Monthly reporting & transparency to investors

7. Next Steps & Deliverables

- Final architectural drawings and site plans (scaled PNG generated in this repo)
- Develop investor pitch (slides generated in this repo)
- Build Excel financial pro-forma for scenario analysis (generated in this repo)
- Begin permitting and construction phase

Attachments generated by script

- `site_layout.png` and `site_layout.svg`
- `financial_model.xlsx`
- `pitch_deck.pptx`

"""
    with open(readme_path, 'w', encoding='utf-8') as f:
        f.write(content.strip() + "\n")
    return readme_path


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # Generate site layout images
    spec = SiteSpec()
    png_path, svg_path = generate_site_layout_images(OUTPUT_DIR, spec)
    print(f"Wrote: {png_path}\nWrote: {svg_path}")

    # Generate financial model
    fin_inputs = FinancialInputs()
    xlsx_path = generate_financial_model(OUTPUT_DIR, fin_inputs)
    print(f"Wrote: {xlsx_path}")

    # Generate pitch deck
    pptx_path = generate_pitch_deck(OUTPUT_DIR)
    print(f"Wrote: {pptx_path}")

    # Generate README
    readme_path = generate_readme(OUTPUT_DIR)
    print(f"Wrote: {readme_path}")


if __name__ == '__main__':
    main()